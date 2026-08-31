import json
import subprocess
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PROJECT_ROOT = Path(__file__).resolve().parents[1]
SCRIPT = PROJECT_ROOT / "scripts/generate_mock_pptx.py"
FIXTURE = PROJECT_ROOT / "examples/mock_slide_deck.json"


def _run_cli(*arguments: str) -> dict:
    completed = subprocess.run(
        [sys.executable, str(SCRIPT), *arguments],
        cwd=PROJECT_ROOT,
        check=True,
        capture_output=True,
        text=True,
    )
    return json.loads(completed.stdout)


def test_mock_pptx_cli_generates_rich_deck_without_runtime(tmp_path):
    result = _run_cli(
        "--input",
        str(FIXTURE),
        "--output-dir",
        str(tmp_path),
    )
    deck_path = Path(result["output_path"])
    deck = Presentation(deck_path)

    assert deck_path.is_file()
    assert result["network_media_enabled"] is False
    assert result["source_slide_count"] == 9
    assert result["exported_slide_count"] == 10
    assert len(deck.slides) == 10
    assert any(shape.has_table for slide in deck.slides for shape in slide.shapes)
    assert any(shape.has_chart for slide in deck.slides for shape in slide.shapes)
    assert any(
        "Câu 1: A" in slide.notes_slide.notes_text_frame.text
        for slide in deck.slides
    )


def test_mock_pptx_cli_maps_fake_url_to_local_image(tmp_path):
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (640, 360), color=(15, 118, 110)).save(image_path)
    payload_path = tmp_path / "payload.json"
    payload_path.write_text(
        json.dumps({
            "lesson_title": "Local media",
            "slides": [{
                "slide_id": "s1",
                "slide_type": "content",
                "title": "Ảnh local",
                "bullets": ["Không gọi mạng"],
                "media": [{"url": "mock://diagram", "caption": "Sơ đồ local"}],
            }],
        }),
        encoding="utf-8",
    )
    media_map_path = tmp_path / "media-map.json"
    media_map_path.write_text(
        json.dumps({"mock://diagram": image_path.name}),
        encoding="utf-8",
    )
    output_dir = tmp_path / "exports"

    result = _run_cli(
        "--input",
        str(payload_path),
        "--media-map",
        str(media_map_path),
        "--output-dir",
        str(output_dir),
    )
    deck = Presentation(result["output_path"])

    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        for shape in deck.slides[0].shapes
    )


def test_mock_pptx_cli_reports_invalid_payload(tmp_path):
    payload_path = tmp_path / "invalid.json"
    payload_path.write_text('{"slides": []}', encoding="utf-8")

    completed = subprocess.run(
        [sys.executable, str(SCRIPT), "--input", str(payload_path)],
        cwd=PROJECT_ROOT,
        capture_output=True,
        text=True,
    )

    assert completed.returncode == 2
    assert "danh sách không rỗng" in completed.stderr
