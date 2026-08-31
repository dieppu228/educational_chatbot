import json
from pathlib import Path

import pytest
from pydantic import ValidationError
from pptx import Presentation

from src.llm.services.slide_export_service import SlideExportService
from src.llm.services.slide_render_planner import SlideRenderPlanner
from src.llm.services.slide_template_registry import (
    TemplateContractError,
    TemplateRegistry,
)
from src.llm.services.slide_merger import SlideQualityGate
from src.schemas.slide_schemas import (
    ChartBlock,
    MergedSlide,
    OutlineSlide,
    ProcessBlock,
    TableBlock,
    normalize_slide_blocks,
)


PROJECT_ROOT = Path(__file__).resolve().parents[1]


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def test_typed_blocks_validate_structural_contracts():
    with pytest.raises(ValidationError):
        TableBlock(columns=["A", "B"], rows=[["only-one"]])
    with pytest.raises(ValidationError):
        ChartBlock(
            chart_type="column",
            categories=["A", "B"],
            series=[{"name": "S", "values": [1]}],
        )
    with pytest.raises(ValidationError):
        ProcessBlock(steps=["Chỉ một bước"])
    with pytest.raises(ValidationError):
        OutlineSlide(slide_id="s1", title="Sai layout", layout_hint="unknown")


def test_all_supported_block_types_parse_through_discriminated_union():
    raw_blocks = [
        {"type": "bullets", "items": ["A"]},
        {"type": "paragraph", "text": "Đoạn văn."},
        {"type": "code", "code": "print(1)", "language": "python"},
        {"type": "table", "columns": ["A"], "rows": [["1"]]},
        {"type": "chart", "categories": ["A"], "series": [{"name": "S", "values": [1]}]},
        {"type": "process", "steps": ["A", "B"]},
        {"type": "comparison", "left_title": "A", "left_items": ["1"], "right_title": "B", "right_items": ["2"]},
        {"type": "callout", "text": "Lưu ý", "tone": "tip"},
    ]

    blocks, _ = normalize_slide_blocks(raw_blocks)
    assert [block.type for block in blocks] == [
        "bullets", "paragraph", "code", "table", "chart", "process", "comparison", "callout"
    ]


def test_invalid_rich_blocks_fall_back_to_legacy_bullets():
    blocks, bullets = normalize_slide_blocks(
        [{
            "type": "chart",
            "chart_type": "pie",
            "categories": ["A", "B"],
            "series": [
                {"name": "S1", "values": [1, 2]},
                {"name": "S2", "values": [3, 4]},
            ],
        }],
        fallback_bullets=["Dữ liệu A", "Dữ liệu B"],
        source_chunk_ids=["c1"],
    )

    assert bullets == ["Dữ liệu A", "Dữ liệu B"]
    assert len(blocks) == 1
    assert blocks[0].type == "bullets"
    assert blocks[0].source_chunk_ids == ["c1"]


def test_normalizer_keeps_valid_blocks_when_one_block_is_invalid():
    blocks, bullets = normalize_slide_blocks(
        [
            {"type": "code", "code": "print('ok')", "source_chunk_ids": ["c1"]},
            {"type": "process", "steps": ["invalid"]},
        ],
        fallback_bullets=["Nội dung dự phòng"],
        source_chunk_ids=["c2"],
    )

    assert [block.type for block in blocks] == ["code", "bullets"]
    assert bullets == ["Nội dung dự phòng"]


def test_render_planner_preserves_all_text_while_splitting():
    bullets = [f"Ý số {index}" for index in range(1, 8)]
    paragraph = " ".join(f"Câu nội dung số {index}." for index in range(1, 80))
    plan = SlideRenderPlanner().plan([
        {
            "slide_id": "s1",
            "slide_type": "content",
            "title": "Nội dung dài",
            "blocks": [
                {"type": "bullets", "items": bullets},
                {"type": "paragraph", "text": paragraph},
            ],
        }
    ])

    rendered_bullets = [
        item
        for slide in plan.slides
        for block in slide.blocks
        for item in block.get("items", [])
    ]
    rendered_text = " ".join(
        block.get("text", "")
        for slide in plan.slides
        for block in slide.blocks
    )
    assert rendered_bullets == bullets
    for index in range(1, 80):
        assert f"Câu nội dung số {index}." in rendered_text
    assert all(
        slide.title == "Nội dung dài" or slide.title.endswith("(tiếp theo)")
        for slide in plan.slides
    )


def test_template_registry_rejects_missing_contract_layout(tmp_path):
    manifest = {
        "template_id": "broken",
        "version": "1.0",
        "layouts": {},
    }
    path = tmp_path / "manifest.json"
    path.write_text(json.dumps(manifest), encoding="utf-8")

    with pytest.raises(TemplateContractError, match="thiếu layouts"):
        TemplateRegistry(path)


def test_bundled_template_has_named_layouts_and_valid_contract():
    template = PROJECT_ROOT / "app/templates/academic_vi_slide_template.pptx"
    manifest = PROJECT_ROOT / "app/templates/academic_vi_slide_template.json"
    deck = Presentation(template)
    registry = TemplateRegistry(manifest)

    registry.validate_presentation(deck)
    assert deck.slide_width / deck.slide_height == pytest.approx(16 / 9, abs=0.02)
    assert set(registry.layouts).issubset({layout.name for layout in deck.slide_layouts})


def test_legacy_template_still_exports_with_compatibility_manifest(tmp_path):
    service = SlideExportService(
        template_path=PROJECT_ROOT / "app/templates/default_slide_template.pptx",
        manifest_path=PROJECT_ROOT / "app/templates/default_slide_template.json",
        export_dir=tmp_path,
    )
    meta = service.export_pptx(
        "Legacy",
        [{"slide_id": "s1", "slide_type": "content", "title": "Nội dung", "bullets": ["Ý chính"]}],
    )

    assert meta["template_id"] == "legacy_default"
    assert (tmp_path / meta["file_id"]).is_file()


def test_template_registry_rejects_version_and_placeholder_type(tmp_path):
    manifest_data = json.loads(
        (PROJECT_ROOT / "app/templates/academic_vi_slide_template.json").read_text(
            encoding="utf-8"
        )
    )
    version_path = tmp_path / "version.json"
    version_path.write_text(
        json.dumps({**manifest_data, "version": "2.0"}),
        encoding="utf-8",
    )
    with pytest.raises(TemplateContractError, match="version"):
        TemplateRegistry(version_path)

    manifest_data["layouts"]["EDU_CONTENT"]["roles"]["body"]["types"] = ["PICTURE"]
    type_path = tmp_path / "type.json"
    type_path.write_text(json.dumps(manifest_data), encoding="utf-8")
    registry = TemplateRegistry(type_path)
    with pytest.raises(TemplateContractError, match="placeholder body"):
        registry.validate_presentation(
            Presentation(PROJECT_ROOT / "app/templates/academic_vi_slide_template.pptx")
        )


def test_corrupt_template_never_publishes_download_file(tmp_path):
    corrupt = tmp_path / "corrupt.pptx"
    corrupt.write_bytes(b"not a PowerPoint package")
    export_dir = tmp_path / "exports"
    service = SlideExportService(
        template_path=corrupt,
        manifest_path=PROJECT_ROOT / "app/templates/academic_vi_slide_template.json",
        export_dir=export_dir,
    )

    with pytest.raises(Exception):
        service.export_pptx("Deck lỗi", [])
    assert not list(export_dir.glob("*.pptx"))


def test_quality_gate_accepts_grounded_rich_blocks():
    slides = [
        MergedSlide(
            slide_id="s1",
            slide_type="title",
            title="Bài học",
            source_chunk_ids=["c1"],
        ),
        MergedSlide(
            slide_id="s2",
            slide_type="content",
            title="Ví dụ",
            blocks=[{"type": "code", "code": "print('ok')", "source_chunk_ids": ["c2"]}],
        ),
        MergedSlide(
            slide_id="s3",
            slide_type="summary",
            title="Tổng kết",
            blocks=[{"type": "bullets", "items": ["Ghi nhớ"], "source_chunk_ids": ["c3"]}],
        ),
    ]

    passed, issues = SlideQualityGate().validate(slides)
    assert passed is True
    assert issues == []


def test_title_and_summary_overflow_preserve_all_text():
    title_items = [f"Phụ đề {index}" for index in range(1, 8)]
    summary_items = [f"Tổng kết {index}" for index in range(1, 9)]
    plan = SlideRenderPlanner().plan([
        {
            "slide_id": "s1",
            "slide_type": "title",
            "title": "Bài học",
            "blocks": [{"type": "bullets", "items": title_items}],
        },
        {
            "slide_id": "s2",
            "slide_type": "summary",
            "title": "Tổng kết",
            "blocks": [{"type": "bullets", "items": summary_items}],
        },
    ])
    rendered = [
        item
        for slide in plan.slides
        for block in slide.blocks
        for item in block.get("items", [])
    ]
    assert rendered == [*title_items, *summary_items]


def test_rich_export_creates_native_objects_and_continuation_slides(tmp_path):
    service = SlideExportService(export_dir=tmp_path)
    slides = [
        {
            "slide_id": "s1",
            "slide_type": "title",
            "title": "Cấu trúc dữ liệu",
            "blocks": [{"type": "bullets", "items": ["Tin học 11"]}],
        },
        {
            "slide_id": "s2",
            "slide_type": "content",
            "title": "Ví dụ code",
            "blocks": [{
                "type": "code",
                "language": "python",
                "code": "for item in values:\n    print(item)",
                "source_chunk_ids": ["c1"],
            }],
        },
        {
            "slide_id": "s3",
            "slide_type": "content",
            "title": "Bảng dữ liệu",
            "blocks": [{
                "type": "table",
                "columns": ["Tên", "Giá trị"],
                "rows": [[f"Mục {idx}", str(idx)] for idx in range(1, 10)],
            }],
        },
        {
            "slide_id": "s4",
            "slide_type": "content",
            "title": "Biểu đồ",
            "blocks": [{
                "type": "chart",
                "chart_type": "column",
                "categories": [f"M{idx}" for idx in range(1, 10)],
                "series": [{"name": "Số lượng", "values": list(range(1, 10))}],
            }],
        },
        {
            "slide_id": "s5",
            "slide_type": "content",
            "title": "Quy trình",
            "blocks": [{"type": "process", "steps": [f"Bước {idx}" for idx in range(1, 8)]}],
        },
        {
            "slide_id": "s6",
            "slide_type": "content",
            "title": "So sánh",
            "blocks": [{
                "type": "comparison",
                "left_title": "Tuần tự",
                "left_items": [f"Trái {idx}" for idx in range(1, 7)],
                "right_title": "Song song",
                "right_items": [f"Phải {idx}" for idx in range(1, 7)],
            }],
        },
    ]

    meta = service.export_pptx("Deck rich content", slides)
    deck = Presentation(str(tmp_path / meta["file_id"]))

    assert meta["template_id"] == "academic_vi"
    assert meta["template_version"] == "1.0"
    assert meta["source_slide_count"] == 6
    assert meta["exported_slide_count"] == 10
    assert len(deck.slides) == 10
    assert sum(shape.has_table for slide in deck.slides for shape in slide.shapes) == 2
    assert sum(shape.has_chart for slide in deck.slides for shape in slide.shapes) == 2
    assert any("for item in values" in _slide_text(slide) for slide in deck.slides)
    assert sum(slide.slide_layout.name == "EDU_COMPARISON" for slide in deck.slides) == 2


def test_missing_media_renders_placeholder_without_failing(tmp_path):
    service = SlideExportService(export_dir=tmp_path)
    meta = service.export_pptx(
        "Media placeholder",
        [{
            "slide_id": "s1",
            "slide_type": "content",
            "title": "Mô hình mạng",
            "blocks": [{"type": "bullets", "items": ["Máy khách", "Máy chủ"]}],
            "media": [{"url": None, "caption": "Sơ đồ client-server", "for_slide_id": "s1"}],
        }],
    )
    deck = Presentation(str(tmp_path / meta["file_id"]))
    assert "Sơ đồ client-server" in _slide_text(deck.slides[0])
    assert deck.slides[0].slide_layout.name == "EDU_CONTENT_MEDIA"


def test_long_caption_is_shortened_with_warning(tmp_path):
    service = SlideExportService(export_dir=tmp_path)
    caption = "Mô tả " * 40
    meta = service.export_pptx(
        "Caption dài",
        [{
            "slide_id": "s1",
            "slide_type": "content",
            "title": "Ảnh minh hoạ",
            "bullets": ["Nội dung"],
            "media": [{"url": None, "caption": caption}],
        }],
    )

    assert any("140 ký tự" in warning for warning in meta["warnings"])
