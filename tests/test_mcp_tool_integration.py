import base64
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.config.config import settings
from src.llm.agents import media_research
from src.llm.agents.media_research import MediaResearchAgent
from src.llm.services.slide_export_service import SlideExportService
from src.llm.services.slide_merger import SlideMerger
from src.schemas.agent_protocol import AgentTask
from src.schemas.slide_schemas import AgentResult
from src.tools.bootstrap import init_tool_layer, reset_tool_layer
from src.tools.implementations.web_search_tool import WebSearchTool
from src.tools.services.tavily_service import TavilyServiceError

_ALLOWED_MEDIA_TYPES = {"image", "gif", "animation", "diagram", "infographic"}


class FakeTavilyService:
    def search(self, query, **kwargs):
        return {
            "results": [
                {
                    "title": "Python loop",
                    "url": "https://example.edu/loop",
                    "content": "Animated for loop",
                    "score": 0.95,
                }
            ],
            "images": [
                {
                    "url": "https://example.edu/assets/loop.gif",
                    "description": "animated for loop",
                }
            ],
        }


class FailingTavilyService:
    def search(self, query, **kwargs):
        raise TavilyServiceError("mock unavailable")


def make_media_task():
    return AgentTask(
        task_id="media-test",
        from_agent="content_supervisor",
        to_agent="media_research_agent",
        task_type="research_media",
        objective="Tìm media minh họa",
        inputs={"topic": "Vòng lặp for", "grade": "10", "book": "KNTT"},
        expected_artifact="media_payload",
    )


def test_tool_layer_registers_system_tools():
    reset_tool_layer()
    client = init_tool_layer()

    names = {tool["name"] for tool in client.list_available_tools()}

    assert names == {"web_search", "content_formatter"}


def test_web_search_returns_gif_through_mcp():
    reset_tool_layer()
    client = init_tool_layer(tavily_service=FakeTavilyService())

    response = client.search_media(
        query="vòng lặp for",
        media_types=["gif"],
        top_k=3,
    )

    assert response.success is True
    assert response.data["media_items"][0]["media_type"] == "gif"
    assert response.data["media_items"][0]["url"].endswith("loop.gif")


def test_web_search_media_parser_skips_invalid_urls_and_accepts_image_url():
    items = WebSearchTool()._parse_media(
        images=[
            {"url": None, "description": None},
            {"image_url": "https://example.edu/assets/loop.png", "description": None},
            {"src": 123, "description": "bad src"},
        ],
        results=[{"url": "https://example.edu/loop", "title": "Loop source"}],
        media_type_filter=["image"],
    )

    assert len(items) == 1
    assert items[0].url == "https://example.edu/assets/loop.png"
    assert items[0].media_type == "image"


def test_web_search_failure_is_returned_without_raise():
    reset_tool_layer()
    client = init_tool_layer(tavily_service=FailingTavilyService())

    response = client.search_media(query="vòng lặp for", media_types=["gif"])

    assert response.success is False
    assert "mock unavailable" in response.error


def test_media_research_agent_enriches_instruction(monkeypatch):
    class FakeMediaAgent:
        def run(self, **kwargs):
            return AgentResult(
                agent="media",
                status="success",
                latency_ms=5,
                payload={
                    "hero_media": [
                        {"caption": "Vòng lặp for chuyển động", "type": "gif", "url": None}
                    ],
                    "inline_media": [],
                },
            )

    class FakeClient:
        def search_media(self, **kwargs):
            assert kwargs["media_types"] == ["gif"]
            return SimpleNamespace(
                success=True,
                data={
                    "media_items": [
                        {
                            "url": "https://example.edu/loop.gif",
                            "media_type": "gif",
                            "source_url": "https://example.edu/loop",
                            "source_title": "Python loop",
                        }
                    ]
                },
            )

    monkeypatch.setattr(media_research, "MediaAgent", FakeMediaAgent)
    monkeypatch.setattr(media_research, "get_mcp_client", lambda: FakeClient())

    result = MediaResearchAgent()._run(make_media_task())
    item = result.artifact["hero_media"][0]

    assert item["url"] == "https://example.edu/loop.gif"
    assert item["media_type"] == "gif"
    assert item["source_url"] == "https://example.edu/loop"
    assert "web_search" in result.used_tools


def test_media_research_agent_keeps_placeholder_on_failure(monkeypatch):
    class FakeMediaAgent:
        def run(self, **kwargs):
            return AgentResult(
                agent="media",
                status="success",
                latency_ms=5,
                payload={
                    "hero_media": [
                        {"caption": "Vòng lặp for chuyển động", "type": "gif", "url": None}
                    ],
                    "inline_media": [],
                },
            )

    class FailingClient:
        def search_media(self, **kwargs):
            return SimpleNamespace(success=False, data=None)

    monkeypatch.setattr(media_research, "MediaAgent", FakeMediaAgent)
    monkeypatch.setattr(media_research, "get_mcp_client", lambda: FailingClient())

    result = MediaResearchAgent()._run(make_media_task())

    assert result.artifact["hero_media"][0]["url"] is None
    assert result.used_tools == ["llm_generation"]


def test_media_enrichment_limits_lookups_and_normalizes_type(monkeypatch):
    calls = []

    class FakeClient:
        def search_media(self, **kwargs):
            calls.append(kwargs)
            return SimpleNamespace(success=False, data=None)

    monkeypatch.setattr(media_research, "get_mcp_client", lambda: FakeClient())
    payload = {
        "hero_media": [
            {"caption": f"Media {index}", "type": "unknown", "url": None}
            for index in range(5)
        ],
        "inline_media": [],
    }

    touched = MediaResearchAgent()._enrich_media_urls(
        payload,
        topic="Tin học",
        grade="10",
        book="KNTT",
    )

    assert touched is False
    assert len(calls) == 4
    assert all(call["media_types"] == ["image"] for call in calls)
    assert payload["hero_media"][0]["type"] == "image"


def test_media_enrichment_parallel_populates_all_items(monkeypatch):
    class FakeClient:
        def search_media(self, **kwargs):
            query = kwargs["query"]
            slug = query.lower().replace(" ", "-")
            return SimpleNamespace(
                success=True,
                data={
                    "media_items": [
                        {
                            "url": f"https://example.edu/{slug}.png",
                            "media_type": kwargs["media_types"][0],
                            "source_url": f"https://example.edu/source/{slug}",
                            "source_title": query,
                        }
                    ]
                },
            )

    monkeypatch.setattr(media_research, "get_mcp_client", lambda: FakeClient())
    payload = {
        "hero_media": [
            {"caption": "Media Alpha", "type": "image", "url": None},
            {"caption": "Media Beta", "type": "gif", "url": None},
        ],
        "inline_media": [
            {"caption": "Media Gamma", "type": "diagram", "url": None},
            {"caption": "Media Delta", "type": "infographic", "url": None},
        ],
    }

    touched = MediaResearchAgent()._enrich_media_urls(
        payload,
        topic="Tin học",
        grade="10",
        book="KNTT",
    )

    all_items = payload["hero_media"] + payload["inline_media"]
    assert touched is True
    assert [item["url"] for item in all_items] == [
        "https://example.edu/media-alpha.png",
        "https://example.edu/media-beta.png",
        "https://example.edu/media-gamma.png",
        "https://example.edu/media-delta.png",
    ]
    assert [item["media_type"] for item in all_items] == [
        "image",
        "gif",
        "diagram",
        "infographic",
    ]


def test_media_enrichment_prioritizes_inline_and_skips_exercise(monkeypatch):
    calls = []

    class FakeClient:
        def search_media(self, **kwargs):
            calls.append(kwargs["query"])
            slug = kwargs["query"].lower().replace(" ", "-")
            return SimpleNamespace(
                success=True,
                data={
                    "media_items": [
                        {
                            "url": f"https://example.edu/{slug}.png",
                            "media_type": kwargs["media_types"][0],
                        }
                    ]
                },
            )

    monkeypatch.setattr(media_research, "get_mcp_client", lambda: FakeClient())
    payload = {
        "hero_media": [
            {"caption": "Hero Alpha", "type": "image", "url": None},
            {"caption": "Hero Beta", "type": "image", "url": None},
        ],
        "inline_media": [
            {"caption": "Content Alpha", "type": "image", "for_slide_type": "content", "url": None},
            {"caption": "Exercise Alpha", "type": "image", "for_slide_type": "exercise", "url": None},
            {"caption": "Content Beta", "type": "image", "for_slide_type": "content", "url": None},
            {"caption": "Content Gamma", "type": "image", "for_slide_type": "content", "url": None},
        ],
    }

    touched = MediaResearchAgent()._enrich_media_urls(
        payload,
        topic="Tin học",
        grade="10",
        book="KNTT",
    )

    assert touched is True
    assert len(calls) == 4
    assert "Exercise Alpha" not in calls
    assert payload["inline_media"][0]["url"]
    assert payload["inline_media"][1]["url"] is None
    assert payload["inline_media"][2]["url"]
    assert payload["inline_media"][3]["url"]


def test_slide_merger_preserves_enriched_media_fields():
    slides = SlideMerger().merge(
        outline_result=AgentResult(
            agent="outline",
            status="success",
            latency_ms=0,
            payload={
                "lesson_title": "Vòng lặp",
                "slides": [
                    {
                        "slide_id": "s1",
                        "slide_type": "title",
                        "title": "Vòng lặp",
                        "key_points": [],
                        "source_chunk_ids": ["c1"],
                    }
                ],
            },
        ),
        content_result=AgentResult(
            agent="content", status="success", latency_ms=0, payload={"slides": []}
        ),
        media_result=AgentResult(
            agent="media",
            status="success",
            latency_ms=0,
            payload={
                "hero_media": [
                    {
                        "caption": "Vòng lặp chuyển động",
                        "type": "animation",
                        "media_type": "animation",
                        "url": "https://example.edu/loop.webm",
                        "source_url": "https://example.edu/loop",
                        "source_title": "Loop animation",
                    }
                ],
                "inline_media": [],
            },
        ),
        quiz_result=AgentResult(
            agent="quiz", status="failed", latency_ms=0, payload={}
        ),
    )

    media = slides[0].media[0]
    assert media.url == "https://example.edu/loop.webm"
    assert media.type == "animation"
    assert media.media_type == "animation"
    assert media.source_url == "https://example.edu/loop"


def test_slide_merger_matches_media_by_slide_id_without_reuse():
    slides = SlideMerger().merge(
        outline_result=AgentResult(
            agent="outline",
            status="success",
            latency_ms=0,
            payload={
                "lesson_title": "Máy tính",
                "slides": [
                    {"slide_id": "s1", "slide_type": "title", "title": "Mở đầu"},
                    {"slide_id": "s2", "slide_type": "content", "title": "CPU"},
                    {"slide_id": "s3", "slide_type": "content", "title": "Internet"},
                    {"slide_id": "s4", "slide_type": "exercise", "title": "Bài tập"},
                ],
            },
        ),
        content_result=AgentResult(
            agent="content",
            status="success",
            latency_ms=0,
            payload={
                "slides": [
                    {"slide_id": "s2", "title": "CPU", "bullets": ["Xử lí nhanh"]},
                    {"slide_id": "s3", "title": "Internet", "bullets": ["Kết nối"]},
                    {"slide_id": "s4", "title": "Bài tập", "bullets": ["Câu hỏi"]},
                ]
            },
        ),
        media_result=AgentResult(
            agent="media",
            status="success",
            latency_ms=0,
            payload={
                "hero_media": [],
                "inline_media": [
                    {
                        "for_slide_id": "s2",
                        "caption": "CPU",
                        "url": "https://example.edu/cpu.png",
                    },
                    {
                        "for_slide_id": "s3",
                        "caption": "Internet",
                        "url": "https://example.edu/internet.png",
                    },
                    {
                        "for_slide_id": "s4",
                        "caption": "Không nên dùng cho bài tập",
                        "url": "https://example.edu/exercise.png",
                    },
                ],
            },
        ),
        quiz_result=AgentResult(agent="quiz", status="failed", latency_ms=0, payload={}),
    )

    assert [media.url for media in slides[1].media] == ["https://example.edu/cpu.png"]
    assert [media.url for media in slides[2].media] == ["https://example.edu/internet.png"]
    assert slides[3].media == []


def test_slide_merger_does_not_reuse_legacy_inline_media_url():
    slides = SlideMerger().merge(
        outline_result=AgentResult(
            agent="outline",
            status="success",
            latency_ms=0,
            payload={
                "lesson_title": "Máy tính",
                "slides": [
                    {"slide_id": "s1", "slide_type": "content", "title": "CPU"},
                    {"slide_id": "s2", "slide_type": "content", "title": "Internet"},
                ],
            },
        ),
        content_result=AgentResult(
            agent="content",
            status="success",
            latency_ms=0,
            payload={"slides": []},
        ),
        media_result=AgentResult(
            agent="media",
            status="success",
            latency_ms=0,
            payload={
                "hero_media": [],
                "inline_media": [
                    {
                        "caption": "Generic",
                        "for_slide_type": "content",
                        "url": "https://example.edu/generic.png",
                    }
                ],
            },
        ),
        quiz_result=AgentResult(agent="quiz", status="failed", latency_ms=0, payload={}),
    )

    assert [media.url for media in slides[0].media] == ["https://example.edu/generic.png"]
    assert slides[1].media == []


def test_slide_export_embeds_downloaded_image(monkeypatch):
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    service = SlideExportService()
    monkeypatch.setattr(service, "_download_media", lambda url: png)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    service._add_media_slot(
        slide,
        {
            "media": [
                {
                    "url": "https://example.edu/image.png",
                    "caption": "Ảnh minh họa",
                }
            ]
        },
    )

    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)


def test_slide_export_blocks_local_media_url():
    service = SlideExportService()

    assert service._is_public_media_url("http://127.0.0.1/private.png") is False
