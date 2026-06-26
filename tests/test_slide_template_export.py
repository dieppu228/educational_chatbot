from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.llm.services.slide_export_service import SlideExportService
from src.llm.services.slide_merger import SlideMerger
from src.schemas.slide_schemas import AgentResult


def test_merger_sets_layout_field_for_content_with_media():
    slides = SlideMerger().merge(
        outline_result=AgentResult(
            agent="outline",
            status="success",
            latency_ms=0,
            payload={
                "lesson_title": "Vòng lặp",
                "slides": [
                    {
                        "slide_id": "s1",
                        "slide_type": "content",
                        "title": "Vòng lặp for",
                        "key_points": ["Cú pháp", "Ví dụ"],
                        "source_chunk_ids": ["c1"],
                    }
                ],
            },
        ),
        content_result=AgentResult(
            agent="content",
            status="success",
            latency_ms=0,
            payload={
                "slides": [
                    {
                        "slide_id": "s1",
                        "title": "Vòng lặp for",
                        "bullets": ["Lặp qua dãy"],
                    }
                ]
            },
        ),
        media_result=AgentResult(
            agent="media",
            status="success",
            latency_ms=0,
            payload={
                "hero_media": [],
                "inline_media": [
                    {
                        "caption": "Minh họa vòng lặp",
                        "for_slide_type": "content",
                        "url": "https://example.edu/loop.png",
                    }
                ],
            },
        ),
        quiz_result=AgentResult(agent="quiz", status="failed", latency_ms=0, payload={}),
    )

    assert slides[0].layout == "content_media"


def _png(width: int, height: int) -> bytes:
    buffer = BytesIO()
    Image.new("RGB", (width, height), color=(24, 124, 98)).save(buffer, format="PNG")
    return buffer.getvalue()


def _has_rendered_picture(slide) -> bool:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if getattr(shape, "image", None) is not None:
            return True
    return False


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text)
    return "\n".join(parts)


def test_export_uses_named_layout(tmp_path):
    service = SlideExportService(export_dir=tmp_path)

    meta = service.export_pptx(
        "Bài kiểm thử",
        [
            {
                "slide_id": "s1",
                "slide_type": "content",
                "layout": "content",
                "title": "Khái niệm",
                "bullets": ["Ý chính"],
            }
        ],
    )

    deck = Presentation(str(tmp_path / meta["file_id"]))
    assert deck.slides[0].slide_layout.name == "Title and Content"


def test_export_picks_media_layout_by_image_aspect_and_inserts_picture(monkeypatch, tmp_path):
    service = SlideExportService(export_dir=tmp_path)
    images = {
        "https://example.edu/wide.png": _png(4, 2),
        "https://example.edu/tall.png": _png(2, 4),
    }
    monkeypatch.setattr(service, "_download_media", lambda url: images[url])

    meta = service.export_pptx(
        "Bài kiểm thử media",
        [
            {
                "slide_id": "s1",
                "slide_type": "image",
                "layout": "content_media",
                "title": "Ảnh ngang",
                "media": [{"url": "https://example.edu/wide.png", "caption": "Ảnh ngang"}],
            },
            {
                "slide_id": "s2",
                "slide_type": "image",
                "layout": "content_media",
                "title": "Ảnh dọc",
                "media": [{"url": "https://example.edu/tall.png", "caption": "Ảnh dọc"}],
            },
        ],
    )

    deck = Presentation(str(tmp_path / meta["file_id"]))
    assert deck.slides[0].slide_layout.name == "Content with Caption"
    assert deck.slides[1].slide_layout.name == "Picture with Caption"
    assert _has_rendered_picture(deck.slides[0])
    assert _has_rendered_picture(deck.slides[1])


def test_export_exercise_slide_includes_answer_without_explanation(tmp_path):
    service = SlideExportService(export_dir=tmp_path)

    meta = service.export_pptx(
        "Bài kiểm thử bài tập",
        [
            {
                "slide_id": "s1",
                "slide_type": "exercise",
                "title": "Bài tập luyện tập",
                "questions": [
                    {
                        "question": "Thiết bị nào dùng để kết nối các mạng?",
                        "options": {
                            "A": "Router",
                            "B": "Bàn phím",
                            "C": "Màn hình",
                            "D": "Máy in",
                        },
                        "correct_answer": "A",
                        "explanation": "Router định tuyến gói tin giữa các mạng.",
                    }
                ],
            }
        ],
    )

    deck = Presentation(str(tmp_path / meta["file_id"]))
    text = _slide_text(deck.slides[0])
    assert "Câu 1. Thiết bị nào dùng để kết nối các mạng?" in text
    assert "A. Router" in text
    assert "Đáp án: A" in text
    assert "Router định tuyến gói tin giữa các mạng." not in text


def test_export_splits_exercise_questions_when_answers_are_shown(tmp_path):
    service = SlideExportService(export_dir=tmp_path)
    questions = [
        {
            "question": f"Câu hỏi kiểm thử số {idx}?",
            "options": {"A": "Một", "B": "Hai", "C": "Ba", "D": "Bốn"},
            "correct_answer": "A",
            "explanation": f"Giải thích câu {idx}.",
        }
        for idx in range(1, 5)
    ]

    meta = service.export_pptx(
        "Bài kiểm thử chia slide",
        [
            {
                "slide_id": "s1",
                "slide_type": "exercise",
                "title": "Bài tập luyện tập",
                "questions": questions,
            }
        ],
    )

    deck = Presentation(str(tmp_path / meta["file_id"]))
    assert len(deck.slides) == 2
    first_text = _slide_text(deck.slides[0])
    second_text = _slide_text(deck.slides[1])
    assert "Bài tập luyện tập (1/2)" in first_text
    assert "Câu hỏi kiểm thử số 1?" in first_text
    assert "Câu hỏi kiểm thử số 3?" not in first_text
    assert "Bài tập luyện tập (2/2)" in second_text
    assert "Câu hỏi kiểm thử số 3?" in second_text
    assert "Đáp án: A" in second_text
