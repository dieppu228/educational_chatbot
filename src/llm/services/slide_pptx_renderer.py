from __future__ import annotations

from io import BytesIO
from typing import Callable, Dict, Iterable, List, Optional

from PIL import Image
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.llm.services.slide_render_models import RenderSlideModel
from src.llm.services.slide_template_registry import TemplateRegistry


CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


class PptxSlideRenderer:
    def __init__(self, registry: TemplateRegistry, media_loader: Callable[[str], Optional[bytes]]):
        self.registry = registry
        self.media_loader = media_loader
        self.warnings: List[str] = []
        theme = registry.info.theme
        self.font = theme.get("font", "Arial")
        self.code_font = theme.get("code_font", "Courier New")
        self.colors = {
            key: RGBColor.from_string(theme.get(key, fallback))
            for key, fallback in {
                "navy": "17324D", "teal": "0F766E", "mint": "DDF4F1",
                "amber": "F59E0B", "background": "F7FAFC", "text": "1F2937",
            }.items()
        }

    def render(self, prs, model: RenderSlideModel, index: int):
        slide = prs.slides.add_slide(self.registry.layout(prs, model.layout_id))
        self._style_background(slide)
        self._set_text(self.registry.placeholder(slide, model.layout_id, "title"), [model.title], 30, bold=True)

        if model.layout_id == "EDU_TITLE":
            self._render_title(slide, model)
        elif model.layout_id == "EDU_COMPARISON":
            self._render_comparison(slide, model)
        elif model.layout_id == "EDU_EXERCISE":
            self._render_exercise(slide, model)
        elif model.layout_id == "EDU_ANSWER_KEY":
            self._render_answer_key(slide, model)
        else:
            self._render_content(slide, model)
        self._write_notes(slide, model)
        return slide

    def _render_title(self, slide, model):
        subtitle = self._block_lines(model.blocks)[:3]
        self._set_text(self.registry.placeholder(slide, model.layout_id, "subtitle"), subtitle, 20)

    def _render_content(self, slide, model):
        block = model.blocks[0] if model.blocks else {"type": "bullets", "items": []}
        kind = block.get("type", "bullets")
        if kind == "code":
            self._render_code(slide, model, block)
        elif kind == "table":
            self._render_table(slide, model, block)
        elif kind == "chart":
            self._render_chart(slide, model, block)
        elif kind == "process":
            self._render_process(slide, model, block)
        elif kind == "callout":
            self._render_callout(slide, model, block)
        else:
            body = self.registry.placeholder(slide, model.layout_id, "body")
            lines = self._block_lines([block])
            self._set_text(
                body,
                lines,
                self._body_font_size(lines),
                bulleted=kind == "bullets",
            )

        if model.media:
            self._render_media(slide, model)

    def _render_code(self, slide, model, block):
        left, top, width, height = self._content_region(slide, model, block)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self.colors["navy"]
        box.line.fill.background()
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Inches(0.24)
        frame.margin_right = Inches(0.24)
        frame.margin_top = Inches(0.18)
        frame.word_wrap = False
        paragraph = frame.paragraphs[0]
        paragraph.text = str(block.get("code") or "")
        paragraph.font.name = self.code_font
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    def _render_table(self, slide, model, block):
        columns = [str(item) for item in block.get("columns") or []]
        rows = [[str(cell) for cell in row] for row in block.get("rows") or []]
        left, top, width, height = self._content_region(slide, model, block)
        graphic = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
        table = graphic.table
        for col, value in enumerate(columns):
            table.cell(0, col).text = value
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row):
                table.cell(row_idx, col_idx).text = value
        column_count = len(columns)
        for cell_idx, cell in enumerate(table.iter_cells()):
            row_idx = cell_idx // column_count
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = self.colors["teal"]
            elif row_idx % 2 == 0:
                cell.fill.fore_color.rgb = self.colors["mint"]
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = self.font
                paragraph.font.size = Pt(14)
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def _render_chart(self, slide, model, block):
        data = CategoryChartData()
        data.categories = [str(item) for item in block.get("categories") or []]
        for series in block.get("series") or []:
            data.add_series(str(series.get("name") or "Dữ liệu"), tuple(series.get("values") or []))
        left, top, width, height = self._content_region(slide, model, block)
        chart = slide.shapes.add_chart(
            CHART_TYPES.get(block.get("chart_type"), XL_CHART_TYPE.COLUMN_CLUSTERED),
            left, top, width, height, data,
        ).chart
        chart.has_legend = len(block.get("series") or []) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.font.name = self.font
        chart.font.size = Pt(12)
        palette = [self.colors["navy"], self.colors["teal"], self.colors["amber"]]
        for index, series in enumerate(chart.series):
            color = palette[index % len(palette)]
            if block.get("chart_type") == "line":
                series.format.line.color.rgb = color
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
        if block.get("chart_type") in {"pie", "doughnut"} and chart.series:
            for index, point in enumerate(chart.series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = palette[index % len(palette)]

    def _render_process(self, slide, model, block):
        steps = [str(item) for item in block.get("steps") or []]
        left, top, width, height = self.registry.region(slide, model.layout_id, "body")
        gap = Inches(0.18)
        box_width = int((width - gap * max(len(steps) - 1, 0)) / max(len(steps), 1))
        box_height = min(height, Inches(1.75))
        y = top + int((height - box_height) / 2)
        for idx, step in enumerate(steps):
            x = left + idx * (box_width + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["mint"]
            shape.line.color.rgb = self.colors["teal"]
            self._set_text(shape, [f"{idx + 1}. {step}"], 16, bold=True, align=PP_ALIGN.CENTER)
            if idx:
                start_x = x - gap
                connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, y + box_height // 2, x, y + box_height // 2)
                connector.line.color.rgb = self.colors["teal"]

    def _render_callout(self, slide, model, block):
        left, top, width, height = self.registry.region(slide, model.layout_id, "body")
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        tone = block.get("tone")
        box.fill.fore_color.rgb = self.colors["amber"] if tone == "warning" else self.colors["mint"]
        box.line.color.rgb = self.colors["teal"]
        lines = [str(block.get("text") or "")]
        self._set_text(box, lines, self._body_font_size(lines), bold=tone == "warning")

    def _render_comparison(self, slide, model):
        block = model.blocks[0] if model.blocks else {}
        fields = [
            ("left_header", [block.get("left_title") or "Phương án A"], 18, True),
            ("left_body", block.get("left_items") or [], 18, False),
            ("right_header", [block.get("right_title") or "Phương án B"], 18, True),
            ("right_body", block.get("right_items") or [], 18, False),
        ]
        for role, lines, size, bold in fields:
            self._set_text(
                self.registry.placeholder(slide, model.layout_id, role),
                lines,
                size,
                bold=bold,
                bulleted=role.endswith("_body"),
            )

    def _render_exercise(self, slide, model):
        lines = []
        for question in model.questions:
            lines.append(f"Câu {question.get('number')}. {question.get('question', '')}")
            options = question.get("options") or {}
            lines.extend(f"{key}. {options[key]}" for key in ("A", "B", "C", "D") if options.get(key))
            lines.append("")
        self._set_text(self.registry.placeholder(slide, model.layout_id, "body"), lines, 18)

    def _render_answer_key(self, slide, model):
        lines = []
        for entry in model.answer_entries:
            text = f"Câu {entry.get('number')}: {entry.get('answer')}"
            if entry.get("explanation"):
                text += f" — {entry['explanation']}"
            lines.append(text)
        self._set_text(self.registry.placeholder(slide, model.layout_id, "body"), lines, 18)

    def _render_media(self, slide, model):
        media = model.media or {}
        data = self.media_loader(str(media.get("url"))) if media.get("url") else None
        left, top, width, height = self.registry.region(slide, model.layout_id, "media")
        caption = str(media.get("caption") or "")
        if caption:
            caption = self._short_caption(caption, model)
            caption_height = Inches(0.42)
            height -= caption_height
            caption_box = slide.shapes.add_textbox(left, top + height, width, caption_height)
            self._set_text(caption_box, [caption], 11, align=PP_ALIGN.CENTER)
        if data:
            try:
                with Image.open(BytesIO(data)) as image:
                    image_ratio = image.width / image.height
                frame_ratio = width / height
                picture = slide.shapes.add_picture(BytesIO(data), left, top, width=width, height=height)
                if image_ratio > frame_ratio:
                    crop = (1 - frame_ratio / image_ratio) / 2
                    picture.crop_left = picture.crop_right = crop
                elif image_ratio < frame_ratio:
                    crop = (1 - image_ratio / frame_ratio) / 2
                    picture.crop_top = picture.crop_bottom = crop
                return
            except Exception:
                self.warnings.append(f"Không chèn được media cho {model.slide_id}")
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self.colors["mint"]
        box.line.color.rgb = self.colors["teal"]
        placeholder_lines = ["◇ HÌNH MINH HOẠ"]
        if not caption:
            placeholder_lines.append("Ảnh minh hoạ sẽ được bổ sung")
        self._set_text(box, placeholder_lines, 14, align=PP_ALIGN.CENTER)

    def _content_region(self, slide, model, block):
        left, top, width, height = self.registry.region(slide, model.layout_id, "body")
        caption = str(block.get("caption") or "")
        if not caption:
            return left, top, width, height
        caption = self._short_caption(caption, model)
        caption_height = Inches(0.38)
        content_height = height - caption_height
        box = slide.shapes.add_textbox(left, top + content_height, width, caption_height)
        self._set_text(box, [caption], 11, align=PP_ALIGN.CENTER)
        return left, top, width, content_height

    def _short_caption(self, caption: str, model: RenderSlideModel) -> str:
        if len(caption) <= 140:
            return caption
        self.warnings.append(f"Caption của {model.slide_id} đã được rút gọn còn 140 ký tự")
        return caption[:137].rstrip() + "..."

    def _style_background(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["background"]

    def _set_text(
        self,
        shape,
        lines: Iterable[str],
        size: int,
        *,
        bold=False,
        align=None,
        bulleted=False,
    ):
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        clean_lines = [str(line) for line in lines if str(line).strip()]
        for index, line in enumerate(clean_lines or [""]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.font.name = self.font
            paragraph.font.size = Pt(size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_after = Pt(5)
            if bulleted:
                paragraph.level = 0
            if align is not None:
                paragraph.alignment = align

    @staticmethod
    def _body_font_size(lines: Iterable[str]) -> int:
        estimated_lines = sum(max(1, (len(str(line)) + 57) // 58) for line in lines)
        if estimated_lines >= 10:
            return 18
        if estimated_lines >= 8:
            return 19
        return 21

    @staticmethod
    def _block_lines(blocks) -> List[str]:
        if isinstance(blocks, dict):
            blocks = [blocks]
        lines = []
        for block in blocks or []:
            kind = block.get("type")
            if kind == "bullets":
                lines.extend(str(item) for item in block.get("items") or [])
            elif kind == "paragraph":
                lines.append(str(block.get("text") or ""))
            elif kind == "callout":
                lines.append(str(block.get("text") or ""))
        return lines

    def _write_notes(self, slide, model):
        parts = [model.notes] if model.notes else []
        if model.source_chunk_ids:
            parts.append("Nguồn chunk: " + ", ".join(model.source_chunk_ids))
        if model.media:
            source = model.media.get("source_url") or model.media.get("url")
            if source:
                parts.append("Nguồn media: " + str(source))
        if parts:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)


__all__ = ["PptxSlideRenderer"]
