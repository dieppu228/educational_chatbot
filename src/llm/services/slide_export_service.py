from __future__ import annotations

import os
import re
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.config.config import project_path, settings
from src.llm.services.slide_media_loader import SlideMediaLoader
from src.llm.services.slide_pptx_renderer import PptxSlideRenderer
from src.llm.services.slide_render_planner import SlideRenderPlanner
from src.llm.services.slide_template_registry import TemplateRegistry


class SlideExportService:
    """Facade for deterministic, template-driven PPTX rendering."""

    def __init__(
        self,
        template_path: Optional[Path] = None,
        export_dir: Optional[Path] = None,
        download_base_url: Optional[str] = None,
        manifest_path: Optional[Path] = None,
    ):
        self.template_path = Path(
            template_path or project_path(settings.SLIDE_TEMPLATE_PATH)
        )
        self.manifest_path = Path(
            manifest_path or project_path(settings.SLIDE_TEMPLATE_MANIFEST_PATH)
        )
        self.export_dir = Path(export_dir or project_path(settings.SLIDE_EXPORT_DIR))
        self.download_base_url = (
            download_base_url or settings.SLIDE_DOWNLOAD_BASE_URL
        ).rstrip("/")
        self.media_loader = SlideMediaLoader()
        self.planner = SlideRenderPlanner()

    def export_pptx(self, lesson_title: str, slides: List[Dict[str, Any]]) -> Dict[str, Any]:
        self.export_dir.mkdir(parents=True, exist_ok=True)
        self.media_loader.clear()

        registry = TemplateRegistry(self.manifest_path)
        prs = Presentation(str(self.template_path))
        registry.validate_presentation(prs)
        prs.core_properties.title = lesson_title or "Bài giảng"
        prs.core_properties.subject = (
            f"EduBot template {registry.info.template_id} v{registry.info.version}"
        )

        plan = self.planner.plan(slides)
        renderer = PptxSlideRenderer(registry, self._download_media)
        for index, model in enumerate(plan.slides, start=1):
            renderer.render(prs, model, index)

        warnings = [*plan.warnings, *renderer.warnings]
        warnings.extend(self._validate_slide_geometry(prs))
        warnings.extend(self._validate_rendered_content(prs, plan.slides))
        file_id = f"{uuid.uuid4().hex}.pptx"
        filename = f"{self._slugify(lesson_title) or 'slide_bai_giang'}.pptx"
        output_path = self.export_dir / file_id
        temp_path = self.export_dir / f".{file_id}.tmp"
        try:
            prs.save(temp_path)
            reopened = Presentation(str(temp_path))
            if len(reopened.slides) != len(plan.slides):
                raise ValueError("PPTX validation failed: slide count mismatch")
            os.replace(temp_path, output_path)
        finally:
            if temp_path.exists():
                temp_path.unlink()

        return {
            "file_id": file_id,
            "filename": filename,
            "download_url": f"{self.download_base_url}/{file_id}",
            "format": "pptx",
            "template_id": registry.info.template_id,
            "template_version": registry.info.version,
            "source_slide_count": plan.source_slide_count,
            "exported_slide_count": len(plan.slides),
            "warnings": warnings,
        }

    def _download_media(self, url: str) -> Optional[bytes]:
        return self.media_loader.download(url)

    @staticmethod
    def _is_public_media_url(url: str) -> bool:
        return SlideMediaLoader.is_public_url(url)

    @staticmethod
    def _validate_slide_geometry(prs) -> List[str]:
        warnings = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.left < 0 or shape.top < 0:
                    warnings.append(f"Slide {slide_index}: shape nằm ngoài canvas")
                    break
                if shape.left + shape.width > prs.slide_width:
                    warnings.append(f"Slide {slide_index}: shape vượt chiều rộng canvas")
                    break
                if shape.top + shape.height > prs.slide_height:
                    warnings.append(f"Slide {slide_index}: shape vượt chiều cao canvas")
                    break
        return warnings

    @staticmethod
    def _validate_rendered_content(prs, models) -> List[str]:
        warnings = []
        for slide_index, (slide, model) in enumerate(zip(prs.slides, models), start=1):
            title_idx = 0
            visible_content = False
            for shape in slide.shapes:
                if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
                    visible_content = True
                    break
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    visible_content = True
                    break
                text = str(getattr(shape, "text", "") or "").strip()
                if not text:
                    continue
                is_title = False
                if getattr(shape, "is_placeholder", False):
                    is_title = shape.placeholder_format.idx == title_idx
                if not is_title:
                    visible_content = True
                    break
            expects_content = bool(
                model.blocks or model.media or model.questions or model.answer_entries
            )
            if expects_content and model.layout_id != "EDU_TITLE" and not visible_content:
                warnings.append(f"Slide {slide_index}: vùng nội dung đang rỗng")
        return warnings

    @staticmethod
    def _slugify(value: str) -> str:
        value = re.sub(r"[^\w\s-]", "", value or "", flags=re.UNICODE).strip().lower()
        return re.sub(r"[-\s]+", "_", value)[:80]


__all__ = ["SlideExportService"]
